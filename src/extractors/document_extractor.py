"""
Document Extractor Module - Enhanced with AUTOMATIC OCR
Extracts text from PPTX, XLSX, and PDF files
NOW WITH: Automatic OCR for scanned PDFs and images
Author: Document Extraction Specialist + Parvesh (OCR Integration)
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import pdfplumber
import json
import os
from pathlib import Path
from typing import Dict, Any, List
from PIL import Image
import io
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# Try to import OCR functionality
try:
    from .ocr_extractor import (
        ocr_pil_image,
        get_pil_image_description,
        extract_images_from_pdf_pages
    )
    OCR_AVAILABLE = True
    print("✅ OCR functionality loaded successfully!")
except ImportError:
    try:
        # Try without relative import (for direct execution)
        import sys
        sys.path.append(os.path.dirname(__file__))
        from ocr_extractor import (
            ocr_pil_image,
            get_pil_image_description,
            extract_images_from_pdf_pages
        )
        OCR_AVAILABLE = True
        print("✅ OCR functionality loaded successfully!")
    except ImportError:
        OCR_AVAILABLE = False
        print("⚠️ OCR not available. Install: pip install pytesseract opencv-python pdf2image")
        print("   For image captioning: pip install transformers torch")


def extract_pptx_with_ocr(file_path: str, use_ocr: bool = True) -> Dict[str, Any]:
    """
    Extract all text from a PowerPoint file, including text in images using OCR.
    
    Args:
        file_path: Path to the .pptx file
        use_ocr: Whether to use OCR for images (default: True)
        
    Returns:
        Dictionary with text content and image information
    """
    try:
        prs = Presentation(file_path)
        text_content = []
        image_texts = []
        image_count = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_text.append(f"=== SLIDE {slide_num} ===")
            
            for shape in slide.shapes:
                # Extract regular text
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        slide_text.append(" | ".join(row_text))
                
                # Extract text from images using OCR
                if use_ocr and OCR_AVAILABLE and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_count += 1
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Convert to PIL Image
                        pil_image = Image.open(io.BytesIO(image_bytes))
                        
                        # OCR the image
                        ocr_text = ocr_pil_image(pil_image)
                        
                        if ocr_text.strip():
                            slide_text.append(f"\n[IMAGE {image_count} - OCR TEXT]")
                            slide_text.append(ocr_text.strip())
                            image_texts.append({
                                "slide": slide_num,
                                "image_num": image_count,
                                "text": ocr_text.strip()
                            })
                        
                        # Get AI description (if available)
                        description = get_pil_image_description(pil_image)
                        if description != "[Image captioning not available]":
                            slide_text.append(f"[IMAGE DESCRIPTION: {description}]")
                    
                    except Exception as e:
                        print(f"  Warning: Could not OCR image on slide {slide_num}: {str(e)}")
            
            if len(slide_text) > 1:  # Has content beyond header
                text_content.append("\n".join(slide_text))
        
        full_text = "\n\n".join(text_content)
        
        return {
            "text": full_text,
            "images_processed": image_count,
            "image_texts": image_texts,
            "ocr_used": use_ocr and OCR_AVAILABLE
        }
    
    except Exception as e:
        return {
            "text": f"Error extracting PPTX: {str(e)}",
            "images_processed": 0,
            "image_texts": [],
            "ocr_used": False
        }


def extract_pptx(file_path: str) -> str:
    """
    Extract all text from a PowerPoint file (backward compatible).
    
    Args:
        file_path: Path to the .pptx file
        
    Returns:
        String containing all extracted text from slides
    """
    result = extract_pptx_with_ocr(file_path, use_ocr=OCR_AVAILABLE)
    return result["text"]


def extract_xlsx(file_path: str) -> str:
    """
    Extract all text from an Excel file.
    
    Args:
        file_path: Path to the .xlsx file
        
    Returns:
        String containing all extracted text from all sheets
    """
    try:
        excel_file = pd.ExcelFile(file_path)
        text_content = []
        
        for sheet_name in excel_file.sheet_names:
            text_content.append(f"=== SHEET: {sheet_name} ===")
            
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Convert DataFrame to readable text
            if not df.empty:
                # Get column headers
                headers = " | ".join(str(col) for col in df.columns)
                text_content.append(headers)
                text_content.append("-" * len(headers))
                
                # Get row data
                for _, row in df.iterrows():
                    row_text = " | ".join(str(val) for val in row.values)
                    text_content.append(row_text)
            else:
                text_content.append("(Empty sheet)")
            
            text_content.append("")  # Add spacing between sheets
        
        return "\n".join(text_content)
    
    except Exception as e:
        return f"Error extracting XLSX: {str(e)}"


def extract_pdf_smart(file_path: str, use_ocr: bool = True) -> Dict[str, Any]:
    """
    Smart PDF extraction: Try digital extraction first, automatically fall back to OCR if needed.
    
    Args:
        file_path: Path to the .pdf file
        use_ocr: Whether to use OCR for scanned pages (default: True)
        
    Returns:
        Dictionary with text content and OCR information
    """
    try:
        # First, try regular text extraction
        text_content = []
        total_pages = 0
        empty_pages = 0
        
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                
                if text and len(text.strip()) > 20:  # Page has digital text
                    text_content.append(f"=== PAGE {page_num} ===\n{text.strip()}")
                else:
                    empty_pages += 1
        
        # Calculate percentage of empty pages
        empty_percentage = (empty_pages / total_pages) * 100 if total_pages > 0 else 0
        
        # If more than 50% of pages are empty/scanned, use full OCR
        if empty_percentage > 50 and use_ocr and OCR_AVAILABLE:
            print(f"  📷 Detected scanned PDF ({empty_pages}/{total_pages} pages need OCR)")
            print(f"  🔄 Running full OCR scan (this may take 1-2 minutes)...")
            
            # Run full OCR
            ocr_result = extract_pdf_full_ocr(file_path)
            
            return {
                "text": ocr_result["text"],
                "total_pages": total_pages,
                "pages_needing_ocr": list(range(1, total_pages + 1)),
                "ocr_used": True,
                "method": "full_ocr"
            }
        else:
            # Use the digital extraction
            full_text = "\n\n".join(text_content)
            
            return {
                "text": full_text,
                "total_pages": total_pages,
                "pages_needing_ocr": [],
                "ocr_used": False,
                "method": "digital"
            }
    
    except Exception as e:
        return {
            "text": f"Error extracting PDF: {str(e)}",
            "total_pages": 0,
            "pages_needing_ocr": [],
            "ocr_used": False,
            "method": "error"
        }


def extract_pdf(file_path: str) -> str:
    """
    Extract text from PDFs with automatic OCR fallback.
    
    Args:
        file_path: Path to the .pdf file
        
    Returns:
        String containing all extracted text from PDF pages
    """
    result = extract_pdf_smart(file_path, use_ocr=OCR_AVAILABLE)
    return result["text"]


def extract_pdf_full_ocr(file_path: str, poppler_path: str = None) -> Dict[str, Any]:
    """
    Extract text from scanned PDF using full OCR on all pages.
    Use this for scanned documents or when regular extraction fails.
    
    Args:
        file_path: Path to the .pdf file
        poppler_path: Path to poppler binaries (Windows only, not needed for Linux)
        
    Returns:
        Dictionary with OCR text and descriptions
    """
    if not OCR_AVAILABLE:
        return {
            "text": "Error: OCR not available. Install required packages.",
            "descriptions": [],
            "pages": 0
        }
    
    try:
        text, descriptions, page_count = extract_images_from_pdf_pages(
            file_path, 
            dpi=300,
            poppler_path=poppler_path
        )
        
        return {
            "text": text,
            "descriptions": descriptions,
            "pages": page_count,
            "method": "full_ocr"
        }
    
    except Exception as e:
        return {
            "text": f"Error in full OCR: {str(e)}",
            "descriptions": [],
            "pages": 0,
            "method": "full_ocr_failed"
        }


def save_extracted_text(text: str, source_file: str, output_dir: str = "data/extracted", 
                       metadata_extra: dict = None) -> str:
    """
    Save extracted text as JSON file.
    
    Args:
        text: Extracted text content
        source_file: Original filename
        output_dir: Directory to save JSON files
        metadata_extra: Additional metadata to include
        
    Returns:
        Path to saved JSON file
    """
    # Create output directory if it doesn't exist
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    # Generate output filename
    base_name = Path(source_file).stem
    extension = Path(source_file).suffix.lower().replace(".", "")
    output_filename = f"{base_name}_{extension}.json"
    output_path = os.path.join(output_dir, output_filename)
    
    # Create JSON structure
    output_data = {
        "text": text,
        "metadata": {
            "source_file": source_file,
            "file_type": extension,
            "character_count": len(text),
            "extracted_successfully": not text.startswith("Error")
        }
    }
    
    # Add extra metadata if provided
    if metadata_extra:
        output_data["metadata"].update(metadata_extra)
    
    # Save to file
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(output_data, f, ensure_ascii=False, indent=2)
    
    return output_path


def process_file(file_path: str, output_dir: str = "data/extracted", use_ocr: bool = True) -> Dict[str, Any]:
    """
    Process a single file and extract text based on file type.
    
    Args:
        file_path: Path to the file to process
        output_dir: Directory to save extracted JSON
        use_ocr: Whether to use OCR for images (default: True)
        
    Returns:
        Dictionary with processing results
    """
    file_extension = Path(file_path).suffix.lower()
    filename = Path(file_path).name
    
    print(f"Processing: {filename}")
    
    # Extract based on file type
    metadata_extra = {}
    
    if file_extension == ".pptx":
        if use_ocr and OCR_AVAILABLE:
            result = extract_pptx_with_ocr(file_path, use_ocr=True)
            text = result["text"]
            metadata_extra = {
                "images_processed": result["images_processed"],
                "ocr_used": result["ocr_used"]
            }
        else:
            text = extract_pptx(file_path)
    
    elif file_extension in [".xlsx", ".xls"]:
        text = extract_xlsx(file_path)
    
    elif file_extension == ".pdf":
        if use_ocr and OCR_AVAILABLE:
            result = extract_pdf_smart(file_path, use_ocr=True)
            text = result["text"]
            metadata_extra = {
                "total_pages": result["total_pages"],
                "pages_needing_ocr": result.get("pages_needing_ocr", []),
                "ocr_used": result["ocr_used"],
                "extraction_method": result.get("method", "unknown")
            }
        else:
            text = extract_pdf(file_path)
    
    else:
        return {
            "filename": filename,
            "status": "skipped",
            "message": f"Unsupported file type: {file_extension}"
        }
    
    # Save extracted text
    output_path = save_extracted_text(text, filename, output_dir, metadata_extra)
    
    return {
        "filename": filename,
        "status": "success" if not text.startswith("Error") else "error",
        "output_path": output_path,
        "character_count": len(text),
        **metadata_extra
    }


def process_all_files(input_dir: str = "data/raw", output_dir: str = "data/extracted", 
                     use_ocr: bool = True) -> list:
    """
    Process all supported files in the input directory.
    
    Args:
        input_dir: Directory containing raw files
        output_dir: Directory to save extracted JSON files
        use_ocr: Whether to use OCR for images (default: True)
        
    Returns:
        List of processing results for each file
    """
    supported_extensions = [".pptx", ".xlsx", ".xls", ".pdf"]
    results = []
    
    # Check if input directory exists
    if not os.path.exists(input_dir):
        print(f"⚠️ Input directory not found: {input_dir}")
        return results
    
    # Get all files in input directory
    files = [f for f in os.listdir(input_dir) if os.path.isfile(os.path.join(input_dir, f))]
    
    if not files:
        print(f"⚠️ No files found in {input_dir}")
        return results
    
    print(f"\n🔍 Found {len(files)} files in {input_dir}")
    if use_ocr and OCR_AVAILABLE:
        print("📷 OCR enabled - will automatically scan images and detect scanned PDFs")
    elif use_ocr and not OCR_AVAILABLE:
        print("⚠️ OCR requested but not available - install dependencies")
    print("=" * 60)
    
    # Process each file
    for filename in files:
        file_path = os.path.join(input_dir, filename)
        file_extension = Path(filename).suffix.lower()
        
        if file_extension in supported_extensions:
            result = process_file(file_path, output_dir, use_ocr=use_ocr)
            results.append(result)
            
            if result["status"] == "success":
                ocr_info = ""
                if result.get('ocr_used'):
                    if result.get('images_processed', 0) > 0:
                        ocr_info = f" (OCR: {result['images_processed']} images)"
                    elif result.get('extraction_method') == 'full_ocr':
                        ocr_info = f" (Full OCR: {result['total_pages']} pages)"
                
                print(f"✅ {filename} → Extracted {result['character_count']} characters{ocr_info}")
            elif result["status"] == "error":
                print(f"❌ {filename} → Extraction failed")
        else:
            print(f"⏭️  {filename} → Skipped (unsupported type)")
    
    print("=" * 60)
    print(f"\n✨ Processing complete! {len([r for r in results if r['status'] == 'success'])} files extracted successfully.")
    
    return results


if __name__ == "__main__":
    # Run extraction on all files in data/raw/
    print("🌈 Starting Document Extraction Process (with automatic OCR)...")
    results = process_all_files(use_ocr=True)
    
    # Print summary
    print("\n📊 EXTRACTION SUMMARY:")
    print(f"Total files processed: {len(results)}")
    print(f"Successful extractions: {len([r for r in results if r['status'] == 'success'])}")
    print(f"Failed extractions: {len([r for r in results if r['status'] == 'error'])}")
    
    # OCR summary
    ocr_files = [r for r in results if r.get('ocr_used')]
    if ocr_files:
        print(f"Files with OCR: {len(ocr_files)}")
        total_images = sum(r.get('images_processed', 0) for r in ocr_files)
        if total_images > 0:
            print(f"Total images processed: {total_images}")
        full_ocr_files = [r for r in ocr_files if r.get('extraction_method') == 'full_ocr']
        if full_ocr_files:
            print(f"Scanned PDFs processed: {len(full_ocr_files)}")
